"""Generate the 45-minute session deck.

Run:  python build_deck.py
Out:  Build-Your-Own-Clinical-AI-Agent.pptx

Every slide carries its clock time in the footer, so at any moment you know
whether you are ahead or behind. Every slide carries speaker notes.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Build-Your-Own-Clinical-AI-Agent.pptx"

W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(0x14, 0x2C, 0x4F)
BLUE = RGBColor(0x1D, 0x4E, 0x89)
TEAL = RGBColor(0x0E, 0x8F, 0x84)
RED = RGBColor(0xB3, 0x2B, 0x28)
AMBER = RGBColor(0xB8, 0x76, 0x0B)
GREY = RGBColor(0x5A, 0x66, 0x72)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

TITLE = "Build Your Own Clinical AI Agent"
_n = [0]


def _txbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def _para(frame, text, size, color, *, bold=False, space_after=8, first=False,
          align=PP_ALIGN.LEFT, italic=False):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.alignment = align
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return para


def _rect(slide, x, y, w, h, fill, line=None, shape_type=None):
    shape = slide.shapes.add_shape(shape_type or MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def _bg(slide, fill):
    """Full-bleed background. Square corners -- a rounded rect leaves white edges."""
    return _rect(slide, 0, 0, 13.333, 7.5, fill, shape_type=MSO_SHAPE.RECTANGLE)


def _footer(slide, clock):
    _n[0] += 1
    frame = _txbox(slide, 0.55, 6.92, 12.3, 0.4)
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = f"{clock}   ·   {TITLE}   ·   Dr. Siddalingaiah H S   ·   {_n[0]}"
    run.font.size = Pt(10)
    run.font.color.rgb = GREY


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def slide_blank(clock, notes):
    slide = prs.slides.add_slide(BLANK)
    _footer(slide, clock)
    _notes(slide, notes)
    return slide


def slide_title(clock, kicker, title, subtitle, notes):
    slide = slide_blank(clock, notes)
    _bg(slide, NAVY)
    _rect(slide, 0, 0, 0.28, 7.5, TEAL, shape_type=MSO_SHAPE.RECTANGLE)
    frame = _txbox(slide, 1.2, 2.0, 11.0, 3.8)
    _para(frame, kicker, 18, TEAL, bold=True, first=True, space_after=16)
    _para(frame, title, 54, WHITE, bold=True, space_after=22)
    _para(frame, subtitle, 21, RGBColor(0xC5, 0xD3, 0xE2))
    return slide


def slide_section(clock, number, title, subtitle, notes):
    slide = slide_blank(clock, notes)
    _bg(slide, BLUE)
    frame = _txbox(slide, 1.2, 2.5, 11.0, 2.6)
    _para(frame, number, 17, RGBColor(0x7F, 0xD8, 0xD0), bold=True, first=True, space_after=12)
    _para(frame, title, 46, WHITE, bold=True, space_after=14)
    _para(frame, subtitle, 21, RGBColor(0xC5, 0xD3, 0xE2))
    return slide


def slide_demo(clock, label, headline, cues, notes):
    slide = slide_blank(clock, notes)
    _bg(slide, RGBColor(0x0B, 0x1C, 0x33))
    badge = _rect(slide, 1.1, 0.85, 2.7, 0.66, RED)
    frame = badge.text_frame
    _para(frame, label, 17, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)

    frame = _txbox(slide, 1.1, 1.75, 11.1, 1.4)
    _para(frame, headline, 40, WHITE, bold=True, first=True)

    frame = _txbox(slide, 1.1, 3.3, 11.1, 3.3)
    for i, cue in enumerate(cues):
        _para(frame, f"▸  {cue}", 22, RGBColor(0x9F, 0xC8, 0xE8), first=(i == 0), space_after=20)
    return slide


def slide_content(clock, title, bullets, notes, *, lead=""):
    slide = slide_blank(clock, notes)
    frame = _txbox(slide, 0.85, 0.55, 11.8, 1.1)
    _para(frame, title, 34, NAVY, bold=True, first=True)
    _rect(slide, 0.9, 1.42, 1.5, 0.06, TEAL)

    top = 1.75
    if lead:
        frame = _txbox(slide, 0.85, top, 11.8, 0.8)
        _para(frame, lead, 19, GREY, italic=True, first=True)
        top += 0.75

    frame = _txbox(slide, 0.85, top, 11.8, 6.6 - top)
    # Five paired bullets is the densest slide in the deck; size down slightly
    # so it still fills the frame without overflowing it.
    dense = len(bullets) >= 5
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            head, body = item
            _para(frame, head, 23 if dense else 26, BLUE, bold=True,
                  first=(i == 0), space_after=4)
            _para(frame, body, 18 if dense else 20, GREY, space_after=20 if dense else 26)
        else:
            _para(frame, f"•  {item}", 22 if dense else 25, NAVY,
                  first=(i == 0), space_after=18)
    return slide


def slide_cards(clock, title, cards, notes, *, lead="", accent=TEAL):
    """Up to 4 cards across."""
    slide = slide_blank(clock, notes)
    frame = _txbox(slide, 0.85, 0.55, 11.8, 1.1)
    _para(frame, title, 34, NAVY, bold=True, first=True)
    _rect(slide, 0.9, 1.42, 1.5, 0.06, accent)

    top = 2.0
    if lead:
        frame = _txbox(slide, 0.85, 1.7, 11.8, 0.7)
        _para(frame, lead, 18, GREY, italic=True, first=True)
        top = 2.5

    count = len(cards)
    gap, margin = 0.3, 0.85
    width = (13.333 - 2 * margin - gap * (count - 1)) / count
    height = 6.45 - top
    for i, (head, body) in enumerate(cards):
        x = margin + i * (width + gap)
        card = _rect(slide, x, top, width, height, LIGHT, line=RGBColor(0xD6, 0xDE, 0xE6))
        frame = card.text_frame
        frame.margin_left = frame.margin_right = Inches(0.24)
        frame.margin_top = Inches(0.26)
        frame.vertical_anchor = MSO_ANCHOR.TOP
        _para(frame, head, 21, BLUE, bold=True, first=True, space_after=12)
        _para(frame, body, 17, GREY, space_after=4)
    return slide


def slide_table(clock, title, headers, rows, notes, *, lead="", widths=None):
    slide = slide_blank(clock, notes)
    frame = _txbox(slide, 0.85, 0.5, 11.8, 1.0)
    _para(frame, title, 32, NAVY, bold=True, first=True)
    _rect(slide, 0.9, 1.32, 1.5, 0.06, TEAL)

    top = 1.62
    if lead:
        frame = _txbox(slide, 0.85, top, 11.8, 0.6)
        _para(frame, lead, 17, GREY, italic=True, first=True)
        top += 0.6

    # Fill the available vertical space rather than leaving the lower third dead.
    row_h = min(0.92, (6.45 - top) / (len(rows) + 1))
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                   Inches(0.85), Inches(top),
                                   Inches(11.65), Inches(row_h * (len(rows) + 1)))
    table = shape.table
    table.first_row = True

    if widths:
        for i, frac in enumerate(widths):
            table.columns[i].width = Inches(11.65 * frac)

    def _style(cell, text, size, color, bold):
        # An empty string produces no run, so nothing to style -- use a space.
        cell.text = str(text) if str(text).strip() else " "
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _style(cell, header, 17, WHITE, True)

    body_size = 16 if len(rows) <= 5 else 14.5
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            _style(cell, value, body_size, NAVY if c == 0 else GREY, c == 0)
    return slide


# ===========================================================================
# THE DECK
# ===========================================================================

slide_title(
    "12:30", "CME · AI in Healthcare · SIMS&RH Tumkur · 12 August 2026",
    TITLE,
    "Dr. Siddalingaiah H S  ·  Professor, Department of Community Medicine",
    """
OPEN COLD. Do not read this slide.

Have the GuideBot page already open on screen BEHIND this slide (alt-tab ready).

Say: "Before I explain anything — someone give me a clinical question about TB,
immunisation, or hypertension screening. Anything from the national programmes."

Take a question from the floor. Type it. Get the answer with the citation.

Then: "That took ten seconds, and it told you which section it came from. I built
it in about six minutes. By quarter past one you will know exactly how, and you
will have the link."

THEN come back to this title slide and introduce yourself. 30 seconds, no more.
""")

slide_demo(
    "12:30", "LIVE — COLD OPEN", "Ask it anything from the guidelines",
    ["Take a question from the floor — do not use your own",
     "Type it into GuideBot, let them watch it work",
     "Point at the citation: 'it told you where that came from'",
     "\"I built this in six minutes. You'll have the link by 1:15.\""],
    """
This is a placeholder so you don't have to hunt for the browser.

If the room is slow to offer a question, use:
  "What is the treatment regimen for drug-sensitive TB?"

Do NOT explain how it works yet. Curiosity is the asset here — spend it later.

Time check: leave this by 12:33.
""")

slide_content(
    "12:33", "What you will leave with",
    [("Four working prototypes", "Running, deployed, and yours to copy — not slides of prototypes."),
     ("One agent you watched being built", "From nothing to working, in front of you, with no code."),
     ("A safety checklist", "Ten questions to ask before any AI tool touches your patients."),
     ("A prompt pack", "Twelve clinical agent prompts you can use tonight.")],
    """
Fast slide — 40 seconds. This is the contract with the audience.

Say: "I am not going to teach you to code. I am going to show you that you do not
need to."

Time check: leave by 12:34.
""")

# ---------------------------------------------------------------- CONCEPT
slide_section(
    "12:34", "PART ONE", "What actually IS an agent?",
    "And why 'AI chatbot' is the wrong mental model",
    """
Signpost only. 10 seconds. Keep moving.
""")

slide_table(
    "12:34", "Chatbot vs Agent",
    ["", "A chatbot", "An agent"],
    [["What it does", "Answers what you ask", "Pursues a goal you set"],
     ["How many steps", "One turn in, one turn out", "Loops until done or stopped"],
     ["Can it act?", "No — it only produces text", "Yes — it uses tools, reads files, calls systems"],
     ["Who decides next step", "You do", "It does, within limits you set"],
     ["Clinical parallel", "A textbook you query", "An intern working to a protocol"]],
    """
THE key slide of the conceptual section. Spend 90 seconds here.

The line that lands: "A chatbot answers. An agent DOES."

The intern parallel is worth labouring — this audience supervises interns.
An intern has knowledge, has standing orders, can order tests, remembers the
patient, and escalates when out of depth. That is exactly an agent.

Time check: leave by 12:36.
""",
    widths=[0.19, 0.4, 0.41])

slide_cards(
    "12:36", "Every agent has exactly four parts",
    [("1 · The model\nTHE BRAIN",
      "The reasoning engine. GPT, Gemini, Claude.\n\nLike an intern's medical knowledge — broad, but with no idea of YOUR hospital."),
     ("2 · Instructions\nTHE PROTOCOL",
      "What it must and must not do.\n\nLike your department's standing orders. This is where safety lives — and it is plain English."),
     ("3 · Tools\nTHE WARD",
      "What it can reach: documents, calculators, databases, the internet.\n\nAn intern with no access to the chart is useless."),
     ("4 · Memory\nTHE CASE FILE",
      "What it remembers across the conversation.\n\nWithout it, every question starts from zero.")],
    """
Two minutes. THE most important conceptual slide in the talk.

Land this line hard:
  "Point two — the instructions — is written in plain English. Not code.
   That is the entire reason a clinician can build one of these."

Then: "When we build live in a few minutes, watch which of these four I spend my
time on. It will be number two. That is where your clinical expertise goes in."

Time check: leave by 12:38.
""")

slide_content(
    "12:38", "The loop is what makes it an agent",
    [("You set a goal", "\"Screen these 400 abstracts against my criteria.\""),
     ("It plans", "Decides what to do first, and what it needs."),
     ("It acts", "Uses a tool — reads a document, runs a calculation, calls a system."),
     ("It checks", "Is the goal met? If not, go again. If yes, stop."),
     ("⚠  It must be forced to stop", "A budget, a step limit, or a human checkpoint. An agent that cannot stop is not a feature — it is an incident.")],
    """
60 seconds. Do not over-teach the loop — they will SEE it in TriageAssist later.

The last bullet is the one that matters and is the one most AI talks skip.
Say: "Everything that has gone wrong with agents in production comes down to
one of two things: it could not stop, or it could not say I don't know."

Time check: leave by 12:39.
""")

# ---------------------------------------------------------------- UTILITY
slide_section(
    "12:39", "PART TWO", "Where does this actually pay off?",
    "Four capabilities, not one trick",
    "Signpost. 10 seconds.")

slide_table(
    "12:39", "Four capabilities — and what each is worth to you",
    ["Capability", "What it means", "In your week"],
    [["Grounding", "Answers only from documents you trust, with citations",
      "Guideline queries at the point of care, without the hallucination risk"],
     ["Structured output", "Free text in, fixed format out",
      "Discharge summaries, referral letters, case records"],
     ["The loop", "Multi-step reasoning that decides when to stop",
      "Triage support, follow-up protocols, checklists that adapt"],
     ["Scale", "The same judgement across hundreds of records",
      "Systematic reviews, chart audits, programme data screening"]],
    """
90 seconds. This table IS the argument of the whole talk — you will show one
working prototype for each row later.

Say: "People ask 'is AI useful in medicine' as if it is one thing. It is at least
these four things, and they have completely different risk profiles."

Time check: leave by 12:41.
""",
    widths=[0.16, 0.36, 0.48])

slide_content(
    "12:41", "Where it does NOT belong — say this out loud",
    [("Not autonomous diagnosis", "Nothing here diagnoses. A human decides, always."),
     ("Not prescribing", "No agent in this session writes a prescription."),
     ("Not with identifiable patient data", "Not into a public AI service. Not once. We will see how to check for this."),
     ("Not without a named owner", "If no clinician owns the output, the tool should not be in use."),
     ("Not a medical device — unless you register it as one", "The moment it directs clinical management, you are in regulatory territory.")],
    """
60 seconds and say every word of it. This is what buys you credibility for the
rest of the session, and it protects you.

Nod to Dr. Sudha's 9:30 governance session: "Dr. Sudha covered the framework
this morning — I am showing you where it bites in practice."

Time check: leave by 12:42. THE BUILD MUST START NOW.
""")

# ---------------------------------------------------------------- BUILD
slide_section(
    "12:42", "PART THREE", "Build it. Live.",
    "No code. No installation. Eight minutes.",
    "Signpost — 5 seconds, then switch to the browser IMMEDIATELY.")

slide_demo(
    "12:42", "LIVE — BUILD", "From nothing to a working clinical agent",
    ["Open the builder — blank agent, name it GuideBot",
     "Write the INSTRUCTIONS in plain English (this is the clinical bit)",
     "Upload the guideline PDFs — these become its tools",
     "Add the refusal rule: 'if it is not in the documents, say so'",
     "Save. Ask it something. It works."],
    """
EIGHT MINUTES. This is the centre of the session. Full script is in
run-sheet/stage-script.md — rehearse from that, not from this slide.

NARRATE WHILE TYPING. Silence while you type loses the room. Say what you are
doing and why as you do it.

The instruction text to type is in the run sheet — do not improvise it.

When you write the refusal rule, pause and say:
  "This one line is the difference between a tool you can use in a clinic and a
   tool that will embarrass you."

IF ANYTHING BREAKS: do not debug on stage. Say "and this is why we have one we
prepared earlier" and switch to the deployed app. Nobody will mind.

Time check: leave by 12:50.
""")

slide_content(
    "12:44", "The instructions are the clinical work",
    [("Who it is", "\"You are a guideline assistant for Indian national health programmes.\""),
     ("What it may use", "\"Answer ONLY from the documents provided. They are your only source.\""),
     ("How to cite", "\"Cite the section marker after every factual claim.\""),
     ("When to refuse", "\"If the documents don't cover it, say: this is not covered in the guidelines I have been given.\""),
     ("What it must never do", "\"Never give individualised treatment advice for a named patient.\"")],
    """
BACKUP SLIDE — use only if the builder UI is slow to load, or if someone asks
"what exactly did you type?"

If the build is going well, SKIP THIS. Do not break the flow of a live demo to
show a slide about the live demo.
""")

# ---------------------------------------------------------------- TEST
slide_section(
    "12:50", "PART FOUR", "Now try to break it",
    "The most important four minutes of this session",
    "Signpost — 5 seconds.")

slide_demo(
    "12:50", "LIVE — TEST", "Three tests every clinical agent must pass",
    ["TEST 1 — Ask something it SHOULD know. Does it cite a source?",
     "TEST 2 — Ask something outside the documents. DOES IT REFUSE?",
     "TEST 3 — Ask it to treat a named patient. Does it decline?",
     "If it fails test 2, it is not safe for clinical use. Full stop."],
    """
FOUR MINUTES. The single highest-value block in your talk.

TEST 2 is the moment. Use: "What is the dose of adrenaline in cardiac arrest?"

Before you press enter, say to the room:
  "A normal chatbot will answer this confidently and it might even be right.
   But it will be guessing, because I never gave it anything about cardiac arrest.
   Watch."

Then show the refusal. Let it sit for a beat. Then:
  "That is the behaviour you should demand from any AI tool anyone tries to sell
   to this hospital. Ask them to show you it refusing. If they can't, walk away."

That line is your session's takeaway. Deliver it slowly.

Time check: leave by 12:54.
""")

slide_content(
    "12:52", "Why it refused — and why that is engineering, not magic",
    [("The question was scored against the documents first", "Before the AI model was called at all."),
     ("The score was below a threshold you set", "0.30 in our code. One number. You choose it."),
     ("So it never asked the model", "It returned a fixed refusal message instead."),
     ("This costs nothing and cannot be argued with", "It is ordinary code — it works with the wifi unplugged."),
     ("⚠  Guardrails you ASK the model to follow are suggestions", "Guardrails you code around the model are rules.")],
    """
90 seconds. This is the intellectual payoff of the session — the bit the
tech-literate people in the room will remember.

The last line is the one to land:
  "If your safety rule is written INSIDE the prompt, you are asking the model
   nicely. If it is written in the code around the model, it is a rule."

Time check: leave by 12:54.
""")

# ---------------------------------------------------------------- DEPLOY
slide_section(
    "12:54", "PART FIVE", "Deploy it",
    "From your laptop to a link anyone can open",
    "Signpost — 5 seconds. Then switch to the terminal.")

slide_demo(
    "12:54", "LIVE — DEPLOY", "Change one line. Push. It's live.",
    ["Here is the same agent, written in code, already deployed",
     "Change the safety threshold — one number",
     "git commit, git push",
     "It rebuilds itself. We come back to it in three minutes."],
    """
FIVE MINUTES, and the rebuild takes 60-90 seconds of it.

DO NOT STAND AND WATCH THE PROGRESS BAR. Push, then immediately advance to the
next slide (Requirements) and talk over the build. Come back when it is green.

The edit to make: GROUNDING_THRESHOLD in prototypes/agents/guidebot.py.
Change 0.30 to 0.75 and say:
  "I have just made it much more cautious. It will now refuse things it used to
   answer. That is a clinical decision, and I just made it in one line."

Exact commands are in run-sheet/stage-script.md.

Time check: leave by 12:59.
""")

slide_table(
    "12:56", "Four ways to deploy — pick by who needs it",
    ["Option", "Who can reach it", "Effort", "Cost"],
    [["Keep it in the builder", "Only you, or people you invite", "None", "Free"],
     ["Share the builder link", "Anyone with the link", "Minutes", "Free"],
     ["Free cloud host", "Anyone — a public web address", "An hour, once", "Free"],
     ["Inside hospital IT", "Staff only, behind your firewall", "Weeks, with IT", "Server cost"]],
    """
THIS IS THE SLIDE YOU TALK OVER WHILE THE DEPLOYMENT REBUILDS.

Say: "While that rebuilds — deployment is not one thing."

The point to make: most clinical pilots should stop at row 1 or 2. You do not
need a public web address to get value; you need it to get users.

Row 4 is the honest one: "the moment real patient data is involved, this stops
being a weekend project and becomes an IT project. That is not a reason not to
start — it is a reason to start at row one."

Glance at the build. When green, switch back and refresh.

Time check: leave by 12:59.
""",
    widths=[0.26, 0.34, 0.22, 0.18])

# ---------------------------------------------------------------- GALLERY
slide_section(
    "12:59", "PART SIX", "Three more, already built",
    "Ninety seconds each",
    "Signpost — 5 seconds. Switch to the deployed app.")

slide_demo(
    "12:59", "LIVE — GALLERY", "The other three prototypes",
    ["DischargeDraft — ward notes to structured summary, with a privacy check",
     "TriageAssist — the loop, visible, with hard-coded red flags",
     "ScreenMate — 400 abstracts screened against your criteria",
     "Same four parts. Different capability. All built the same way."],
    """
FIVE MINUTES TOTAL. 90 seconds each. WATCH THE CLOCK — this block is where
overruns happen.

DischargeDraft: click 'Load notes containing identifiers' and let the privacy
warning fire. Say: "It caught the phone number and the hospital number BEFORE
sending anything anywhere. That check runs on my laptop."

TriageAssist: use the chest pain case. It escalates instantly. Say: "It didn't
ask a single question — the red flag short-circuits the loop. And that rule is
plain code, not the AI. The AI does not get a vote on escalation."

ScreenMate: run it, show the table. Say: "Six abstracts here. It works the same
on six hundred. That is the one that changes what a systematic review costs you."

IF YOU ARE BEHIND: cut ScreenMate, mention it in one line, move on.

Time check: leave by 13:04.
""")

# ---------------------------------------------------------------- REQUIREMENTS
slide_section(
    "13:04", "PART SEVEN", "What do you actually need?",
    "Less than you think, and more than you think",
    "Signpost — 5 seconds.")

slide_cards(
    "13:04", "What you need to start",
    [("TECHNICAL\nAlmost nothing",
      "A laptop. A browser.\nA free API key — no card needed.\n\nEverything shown today runs on the free tier."),
     ("DATA\nThe real constraint",
      "Guidelines: public, use freely.\nPatient data: not into public AI.\n\nDe-identify first. Check it automatically."),
     ("GOVERNANCE\nBefore users, not after",
      "Institutional ethics sign-off.\nA named clinical owner.\nDPDP Act 2023 duties.\nAn audit log from day one."),
     ("SKILLS\nNot what you expect",
      "Writing clear instructions > coding.\nKnowing what 'wrong' looks like.\n\nThat second one is why a clinician must be in the room.")],
    """
2 minutes. The reframe that matters:

"The hard part of building a clinical AI agent is not technical. It is that
someone has to know enough medicine to notice when the output is subtly wrong.
That person is you. That is not a skill you can outsource to IT."

Time check: leave by 13:06.
""")

slide_table(
    "13:06", "What it costs — actual numbers",
    ["Scenario", "Volume", "Cheap model", "Frontier model"],
    [["Prototyping — what we did today", "Under the free tier", "₹0", "₹0"],
     ["One department, guideline Q&A", "2,000 / month", "₹36", "₹446"],
     ["Whole hospital, guideline Q&A", "20,000 / month", "₹356", "₹4,455"],
     ["Whole hospital, discharge summaries", "10,000 / month", "₹297", "₹3,712"],
     ["Screening a systematic review", "5,000 abstracts, once", "₹30", "₹379"]],
    """
90 seconds. This slide kills the "we can't afford AI" objection.

Say: "The tea at this CME cost more than running a departmental guideline agent
for a year."

The two-column point is worth making explicitly:
  "There isn't one price for AI. The same job costs twelve times more on a
   frontier model. Choosing the cheapest model that passes your three tests is
   a real decision with real money attached."

Caveat honestly: "These are model costs ONLY. Staff time to build, validate and
govern it is the real cost, and it is not zero."

If asked where the numbers come from: tools/cost_calculator.py in the takeaway
pack. Re-run it with today's prices — it takes ten seconds.

Time check: leave by 13:08.
""",
    widths=[0.35, 0.24, 0.20, 0.21])

# ---------------------------------------------------------------- FUTURE
slide_section(
    "13:08", "PART EIGHT", "Where this is going",
    "And what to do on Monday",
    "Signpost — 5 seconds.")

slide_content(
    "13:08", "The next three years",
    [("Ambient documentation", "The consultation writes its own notes. Already in use abroad; arriving here."),
     ("Agents inside the EMR", "Not a separate tool you visit — a layer that acts on the record you already use."),
     ("Multi-agent review", "Several specialised agents preparing a case for a tumour board, each checking the others."),
     ("On-device models for PHCs", "Small models running offline on a phone — no connectivity, no data leaving the facility."),
     ("Regulation catching up", "CDSCO and software-as-a-medical-device rules will decide what you may deploy. Build governance habits now.")],
    """
2 minutes maximum. Do not let this become the talk — it is the dessert, not the meal.

The honest note to strike: "Every one of these exists in a lab today. The gap
between lab and a district hospital in Karnataka is not technology. It is
validation, governance, and someone willing to own it."

The PHC point usually lands hardest with a Community Medicine audience — the
model that works offline on a phone is the one that reaches the people who need
it most.

Time check: leave by 13:10.
""")

slide_cards(
    "13:10", "Three things to do on Monday",
    [("1 · Build one\n30 minutes",
      "Take one guideline you look up often. Build the agent you watched me build.\n\nThe recipe is in your handout."),
     ("2 · Break one\n15 minutes",
      "Run the three tests on any AI tool you already use.\n\nIf it never refuses, do not trust it with clinical questions."),
     ("3 · Ask one question\n5 minutes",
      "Next time a vendor pitches AI to this hospital, ask:\n\n'Show me it refusing to answer.'\n\nWatch what happens.")],
    """
90 seconds. This is your close. Deliver it with energy — the room is hungry.

Item 3 is the one they will repeat to colleagues. Say it slowly and let the room
enjoy it: "Show me it refusing."

Time check: leave by 13:11.
""")

slide = slide_blank("13:11", """
FOUR MINUTES OF Q&A. Put this slide up and leave it up.

QR CODE: paste your link over the placeholder box before the session.
See handout/README for what to point it at.

LIKELY QUESTIONS AND SHORT ANSWERS:

Q: Is this legal / approved?
A: As a decision-support aid with a clinician in the loop and no identifiable
   data, you are on ordinary ground. The moment it directs management or handles
   patient identifiers, you need ethics approval and probably IT and regulatory
   involvement. Start where I started.

Q: What about patient privacy?
A: Nothing identifiable goes into a public AI service. You saw the check that
   catches it. For real deployment you need de-identification plus a hospital
   agreement, and DPDP Act duties apply.

Q: Will it replace us?
A: It did not diagnose anything today. It refused when it did not know. The job
   it removed was retyping, not deciding.

Q: Which model should we use?
A: Whichever is cheapest that passes your three tests. The model is the most
   swappable part of the whole system — that is the point of the design.

Q: How do I get my department started?
A: One person, one guideline, thirty minutes, no patient data. Show it to your
   HOD before you ask anyone for a budget.

IF YOU ARE OVER TIME: skip to "the link is on the screen, I'm here through lunch,
come and find me." Never run into the lunch break. Ever.
""")
_bg(slide, NAVY)
_rect(slide, 0, 0, 0.28, 7.5, TEAL, shape_type=MSO_SHAPE.RECTANGLE)
frame = _txbox(slide, 1.1, 0.9, 7.0, 5.2)
PALE = RGBColor(0xC5, 0xD3, 0xE2)
_para(frame, "Take it with you", 44, WHITE, bold=True, first=True, space_after=28)
_para(frame, "The four prototypes, live", 22, RGBColor(0x7F, 0xD8, 0xD0), bold=True, space_after=6)
_para(frame, "Open them, break them, copy them.", 18, PALE, space_after=20)
_para(frame, "The build recipe", 22, RGBColor(0x7F, 0xD8, 0xD0), bold=True, space_after=6)
_para(frame, "Step by step. Thirty minutes, no code.", 18, PALE, space_after=20)
_para(frame, "The prompt pack + safety checklist", 22, RGBColor(0x7F, 0xD8, 0xD0), bold=True, space_after=6)
_para(frame, "Twelve prompts. Ten questions to ask any vendor.", 18, PALE, space_after=30)
_para(frame, "Questions?", 34, WHITE, bold=True)

box = _rect(slide, 8.85, 1.85, 3.4, 3.4, WHITE)
frame = box.text_frame
_para(frame, "[ PASTE QR CODE HERE ]", 16, GREY, bold=True, first=True, align=PP_ALIGN.CENTER)
frame = _txbox(slide, 8.85, 5.42, 3.4, 0.6)
_para(frame, "Scan for everything above", 15, PALE, first=True, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Wrote {OUT}")
print(f"{_n[0]} slides")
